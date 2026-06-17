from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("outputs/home-ecommerce-material-tagging.pptx")

W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(31, 34, 37)
PAPER = RGBColor(247, 244, 237)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(95, 101, 105)
SAGE = RGBColor(91, 124, 105)
CLAY = RGBColor(196, 91, 58)
OCHRE = RGBColor(218, 169, 87)
MIST = RGBColor(229, 234, 229)
LINE = RGBColor(216, 220, 214)

TITLE_FONT = "Microsoft YaHei UI"
BODY_FONT = "Microsoft YaHei"


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def set_text(box, text, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def text(slide, x, y, w, h, content, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(box, content, size=size, color=color, bold=bold, align=align)


def title(slide, content, subtitle=None, dark=False):
    color = WHITE if dark else INK
    text(slide, 0.68, 0.48, 8.8, 0.62, content, 28, color, True)
    if subtitle:
        text(slide, 0.72, 1.1, 8.8, 0.36, subtitle, 11, WHITE if dark else MUTED)


def pill(slide, x, y, w, label, fill=MIST, color=INK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    set_fill(shape, fill)
    shape.line.color.rgb = fill
    set_text(shape, label, size=10.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def card(slide, x, y, w, h, heading, body, accent=SAGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, WHITE)
    shape.line.color.rgb = LINE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(bar, accent)
    text(slide, x + 0.25, y + 0.18, w - 0.45, 0.32, heading, 15, INK, True)
    text(slide, x + 0.25, y + 0.58, w - 0.45, h - 0.72, body, 10.5, MUTED)
    return shape


def bg(slide, color=PAPER):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(rect, color)
    rect.line.color.rgb = color


def add_footer(slide, n, dark=False):
    color = RGBColor(196, 201, 199) if dark else MUTED
    text(slide, 11.05, 7.0, 1.35, 0.25, f"{n:02d} / 10", 8.5, color, align=PP_ALIGN.RIGHT)


def connector(slide, x1, y1, x2, y2, color=CLAY):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def create_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "家居电商混剪素材打标", "第一版只做三件事：是什么 / 有什么用 / 适合放哪", True)
    text(s, 0.74, 2.15, 7.3, 1.3, "不要让人面对几十个标签。\n让每个素材片段只回答 3 个问题。", 31, WHITE, True)
    for i, label in enumerate(["画面类型", "卖点方向", "剪辑位置"]):
        pill(s, 0.78 + i * 1.55, 4.1, 1.18, label, CLAY if i == 0 else (SAGE if i == 1 else OCHRE), WHITE)
    text(s, 0.78, 6.75, 3.5, 0.28, "Creative Studio / 家居类目混剪工具", 9.5, RGBColor(196, 201, 199))

    # 2
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "为什么会头晕", "问题不在标签本身，而在标签没有服务剪辑决策")
    card(s, 0.75, 1.75, 3.75, 3.2, "复杂打标的问题", "产品类目、镜头语言、场景、人群、质量、风格全摊开，用户会变成素材管理员。", CLAY)
    card(s, 4.8, 1.75, 3.75, 3.2, "混剪真正需要", "系统只需要知道：这个片段能不能放进某个脚本节点或模板槽位。", SAGE)
    card(s, 8.85, 1.75, 3.75, 3.2, "第一版原则", "固定字段少一点，自动建议多一点，人工只确认和改错。", OCHRE)
    text(s, 1.0, 5.75, 10.8, 0.55, "结论：标签体系要围绕“剪辑位置”收敛，不围绕“描述世界”发散。", 20, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "三问法", "每个素材片段只需要被确认一次")
    labels = [
        ("1", "它是什么？", "画面类型", "产品全貌、产品细节、使用动作、生活场景"),
        ("2", "它有什么用？", "卖点方向", "颜值、省空间、收纳、舒适、材质"),
        ("3", "它适合放哪？", "剪辑位置", "开头钩子、产品亮相、卖点展示、结尾 CTA"),
    ]
    for i, (num, h1, h2, body) in enumerate(labels):
        x = 0.88 + i * 4.05
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.9), Inches(0.7), Inches(0.7))
        set_fill(circle, [CLAY, SAGE, OCHRE][i])
        set_text(circle, num, 20, WHITE, True, PP_ALIGN.CENTER)
        text(s, x, 2.82, 3.35, 0.34, h1, 17, INK, True)
        text(s, x, 3.25, 3.35, 0.36, h2, 24, [CLAY, SAGE, OCHRE][i], True)
        text(s, x, 3.92, 3.2, 0.85, body, 12.5, MUTED)
        if i < 2:
            connector(s, x + 3.25, 2.25, x + 3.78, 2.25, LINE)
    text(s, 1.05, 5.85, 11.3, 0.55, "多类目兼容靠这三问，而不是给每个家具品类都设计一套标签。", 19, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "标签组一：画面类型", "描述“画面本身是什么”，用于判断素材基础角色")
    tags = ["产品全貌", "产品细节", "使用动作", "生活场景", "空间氛围", "安装/开箱", "前后对比", "人物出镜", "已剪成片"]
    positions = [(0.9, 1.65), (4.2, 1.65), (7.5, 1.65), (0.9, 2.75), (4.2, 2.75), (7.5, 2.75), (0.9, 3.85), (4.2, 3.85), (7.5, 3.85)]
    for idx, tag in enumerate(tags):
        x, y = positions[idx]
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.72))
        set_fill(shape, WHITE)
        shape.line.color.rgb = LINE
        set_text(shape, tag, 16, INK, True, PP_ALIGN.CENTER)
    text(s, 1.0, 5.7, 10.9, 0.5, "例：沙发、柜子、灯具都可以是“产品全貌”；抽屉开合、坐垫回弹都可以是“使用动作”。", 15, MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "标签组二：卖点方向", "描述“这个片段能证明什么”，用于匹配脚本文案")
    tags = ["颜值", "省空间", "收纳", "舒适", "稳固", "易清洁", "材质", "安装简单", "适合小户型", "适合家庭", "适合养宠", "价格活动"]
    for i, tag in enumerate(tags):
        row, col = divmod(i, 4)
        pill(s, 1.0 + col * 2.9, 1.65 + row * 0.9, 2.15, tag, MIST if i % 3 else RGBColor(235, 226, 215), INK)
    card(s, 1.05, 5.0, 5.15, 1.25, "关键取舍", "卖点标签保持通用，不跟品类绑死。品类差异先放到产品信息里，后续再扩展品类专属卖点。", SAGE)
    card(s, 6.75, 5.0, 5.15, 1.25, "系统用法", "脚本说“省空间”，系统优先找带有省空间、收纳、小户型的片段。", CLAY)
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "标签组三：剪辑位置", "这是混剪模板最关心的一组标签")
    flow = [
        ("开头钩子", CLAY),
        ("产品亮相", SAGE),
        ("场景代入", OCHRE),
        ("卖点展示", SAGE),
        ("细节证明", CLAY),
        ("使用演示", SAGE),
        ("转场过渡", OCHRE),
        ("结尾 CTA", CLAY),
        ("封面候选", INK),
    ]
    for i, (label, color) in enumerate(flow):
        x = 0.65 + (i % 3) * 4.1
        y = 1.55 + (i // 3) * 1.35
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.05), Inches(0.78))
        set_fill(shape, color)
        set_text(shape, label, 16, WHITE, True, PP_ALIGN.CENTER)
    text(s, 0.9, 5.85, 11.5, 0.45, "第一版自动混剪时，先按“剪辑位置”取素材，再用“画面类型”和“卖点方向”排序。", 17, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "一个素材卡片应该长什么样", "用户看到的是选择题，不是标签海", True)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.65), Inches(5.35), Inches(4.8))
    set_fill(panel, WHITE)
    text(s, 1.45, 1.98, 2.2, 0.3, "00:12 - 00:17", 12, MUTED, True)
    thumb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.45), Inches(2.45), Inches(2.25), Inches(2.65))
    set_fill(thumb, RGBColor(225, 229, 224))
    text(s, 1.78, 3.45, 1.6, 0.45, "视频片段预览", 13, MUTED, True, PP_ALIGN.CENTER)
    text(s, 4.0, 2.48, 1.9, 0.28, "画面类型", 10.5, MUTED, True)
    pill(s, 4.0, 2.85, 1.35, "使用动作", SAGE, WHITE)
    text(s, 4.0, 3.38, 1.9, 0.28, "卖点方向", 10.5, MUTED, True)
    pill(s, 4.0, 3.75, 1.05, "收纳", CLAY, WHITE)
    pill(s, 5.15, 3.75, 1.05, "省空间", CLAY, WHITE)
    text(s, 4.0, 4.28, 1.9, 0.28, "剪辑位置", 10.5, MUTED, True)
    pill(s, 4.0, 4.65, 1.55, "卖点展示", OCHRE, WHITE)
    text(s, 7.1, 2.35, 4.8, 1.6, "人工只确认：\n画面类型 / 卖点方向 / 剪辑位置 / 是否可用", 27, WHITE, True)
    text(s, 7.15, 4.55, 4.5, 0.76, "清晰度、比例、水印、光线、是否有人声等信息，优先交给系统自动识别，后台保存。", 14, RGBColor(212, 218, 214))
    add_footer(s, 7, dark=True)

    # 8
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "脚本如何对应素材", "脚本节点不直接找文件，而是找“标签组合”")
    nodes = [
        ("痛点", "小户型 / 生活场景 / 开头钩子"),
        ("产品亮相", "产品全貌 / 颜值 / 产品亮相"),
        ("功能演示", "使用动作 / 收纳 / 卖点展示"),
        ("细节证明", "产品细节 / 稳固或材质 / 细节证明"),
    ]
    for i, (h, b) in enumerate(nodes):
        x = 0.75 + i * 3.05
        card(s, x, 1.75, 2.55, 2.2, h, b, [CLAY, SAGE, OCHRE, CLAY][i])
        if i < 3:
            connector(s, x + 2.55, 2.85, x + 2.9, 2.85, LINE)
    text(s, 1.1, 5.25, 10.9, 0.7, "这样脚本、素材、模板三者就能对齐：脚本提出需求，模板定义槽位，标签负责匹配。", 20, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "第一版工作流", "半自动，而不是一上来全自动")
    steps = [
        ("上传素材", "原片 / AI 视频 / 成片"),
        ("自动切片", "按镜头或手动截取"),
        ("AI 建议标签", "三问法自动预填"),
        ("人工确认", "只修错，不重填"),
        ("模板混剪", "按槽位自动取片段"),
        ("统一包装", "字幕 / 配音 / 封面 / CTA"),
    ]
    for i, (h, b) in enumerate(steps):
        x = 0.65 + (i % 3) * 4.1
        y = 1.55 + (i // 3) * 1.75
        card(s, x, y, 3.15, 1.18, h, b, [CLAY, SAGE, OCHRE][i % 3])
    text(s, 1.0, 5.95, 11.2, 0.45, "用户的真实体验：系统先给一版，用户只改明显不对的素材。", 18, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "建议的落地顺序", "先让素材能被模板调用，再谈更复杂的识别和智能剪辑", True)
    card(s, 0.9, 1.55, 3.65, 3.65, "第 1 步", "先实现素材片段卡片：画面类型、卖点方向、剪辑位置、是否可用。", CLAY)
    card(s, 4.85, 1.55, 3.65, 3.65, "第 2 步", "做一个家居通用模板：开头钩子、产品亮相、卖点展示、细节证明、结尾 CTA。", SAGE)
    card(s, 8.8, 1.55, 3.65, 3.65, "第 3 步", "接入 HyperFrames 做字幕、封面、片头片尾和统一包装。", OCHRE)
    text(s, 1.0, 6.25, 11.3, 0.42, "最小闭环：少量标签 + 一个模板 + 一次可预览的混剪结果。", 21, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s, 10, dark=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


if __name__ == "__main__":
    create_deck()
    print(OUT.resolve())
